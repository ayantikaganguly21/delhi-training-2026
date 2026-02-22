import os
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_docx():
    doc = Document()
    doc.add_heading('Prompt Engineering Documentation', 0)
    
    doc.add_heading('App Objective', level=1)
    doc.add_paragraph('RailFind AI aims to solve the problem of lost & found items on Indian Railways (e.g., HWH to NDLS), by providing an AI-powered matching platform for passengers, TTEs, and Station Masters.')
    
    doc.add_heading('Prompts Used to Design & Generate the App', level=1)
    
    doc.add_heading('Prompt 1: Ideation & Problem Identification', level=2)
    doc.add_paragraph('Prompt: "Act as a product manager. Identify the exact pain points passengers face when losing items on a train journey from Howrah (HWH) to New Delhi (NDLS). Suggest an AI-driven solution."\n\nResult: Identified that lack of centralized tracking, shifting train locations, and poor communication between finders and losers are core issues. Suggested an AI matching engine.')
    
    doc.add_heading('Prompt 2: App UI/UX Design', level=2)
    doc.add_paragraph('Prompt: "Design a responsive frontend structure using HTML, CSS, and JS for RailFind AI. Include a dark glassmorphism aesthetic suitable for a premium tech product. Include views for Report Lost, Report Found, and AI Matches."\n\nResult: Generated a beautifully styled app with glass-card layouts, neon gradients, and a simulated feed of AI matches with probability scores.')
    
    doc.add_heading('Prompt 3: AI Feature Mocking', level=2)
    doc.add_paragraph('Prompt: "Write JavaScript logic to handle form submissions and simulate AI matching. Show an example of an 85% match for a trolley bag found in the B4 coach of the HWH-NDLS Rajdhani."\n\nResult: Implemented the app.js logic with hardcoded mock responses showing exactly how the AI would map descriptions and train details.')
    
    doc.add_heading('How Prompt Refinement Improved Output', level=1)
    doc.add_paragraph('Initially, generic prompts yielded a basic form application. By refining the prompts to specify "dark glassmorphism", "HWH to NDLS specific journey references", and "AI Match probability scores", the resulting application became highly contextualized and visually impressive. It transformed from a simple CRUD app to an intelligent AI dashboard simulation.')

    doc.save('Prompt_Engineering_Documentation.docx')

def create_pptx():
    prs = Presentation()
    
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Lost & Found Inside Train Coach"
    subtitle.text = "AI-Powered Problem Solving for Train Journeys\nBy Antigravity"
    
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Problem Identification"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "The Pain of Losing Items on a Train (HWH to NDLS)"
    p = tf.add_paragraph()
    p.text = "Passengers frequently lose baggage, electronics, or documents."
    p = tf.add_paragraph()
    p.text = "No centralized system to report lost items in real-time."
    p = tf.add_paragraph()
    p.text = "Finders (passengers or TTEs) don't know who the item belongs to."
    p = tf.add_paragraph()
    p.text = "Once the train reaches NDLS, finding items left at intermediate stations is nearly impossible."

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Why We Need Technology"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "Manual reporting is slow and fragmented across different stations."
    p = tf.add_paragraph()
    p.text = "An AI platform can instantly cross-reference 'lost' descriptions with 'found' descriptions."
    p = tf.add_paragraph()
    p.text = "Computer Vision can auto-tag uploaded images of found items."
    p = tf.add_paragraph()
    p.text = "A mobile-first web app allows instant reporting right from the coach seat."

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Introducing RailFind AI"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "Core Features:"
    p = tf.add_paragraph()
    p.text = "Report Lost Item: Input train details, description, and predictive AI mapping."
    p = tf.add_paragraph()
    p.text = "Report Found Item: Auto-categorization of found items via image."
    p = tf.add_paragraph()
    p.text = "AI Matcher: Instantly shows match probability (e.g., 94% Match for a black trolley in B4 coach)."

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Conclusion"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "RailFind AI bridge the communication gap between passengers."
    p = tf.add_paragraph()
    p.text = "AI eliminates the manual effort of searching through logs."
    p = tf.add_paragraph()
    p.text = "Makes train journeys stress-free."

    prs.save('Train_Travel_Presentation.pptx')

if __name__ == '__main__':
    create_docx()
    create_pptx()
    print("Successfully generated Word and PPT files.")
